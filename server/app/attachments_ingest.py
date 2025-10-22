# Project:Confluence Evidence API  Component:attachments_ingest  Version:v1.4.3
from __future__ import annotations
import os
import json
import signal
import hashlib
import sqlite3
import logging
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple

from .config import settings
from .repository import connect
from .chunker_rag import iter_chunks
from .rag_store import NumpyStore, VSConfig
from .embeddings import get_embedder

CHUNK_MIN_CHARS   = int(os.getenv("CHUNK_MIN_CHARS", "50"))
ATT_MAX_CHARS     = int(os.getenv("ATT_MAX_CHARS",   "200000"))
OCR_TIMEOUT_SEC   = int(os.getenv("OCR_TIMEOUT_SEC", "180"))
PDF_TXT_TIMEOUT   = int(os.getenv("PDF_TXT_TIMEOUT", "45"))
IMG_OCR_LANG      = os.getenv("OCR_LANG", "eng")
OCR_DPI           = int(os.getenv("OCR_DPI", "50"))
OCR_ENABLED       = os.getenv("OCR_ENABLED", "0") == "1"
OCR_PDF_IF_EMPTY  = os.getenv("OCR_PDF_IF_EMPTY", "0") == "1"
OCR_IMAGES        = os.getenv("OCR_IMAGES", "0") == "1"
OCR_JOBS          = int(os.getenv("OCR_JOBS", "0"))
FLUSH_EVERY       = int(os.getenv("FLUSH_EVERY", "0"))  # 0 => use batch

ATT_TEXT_CACHE    = Path(os.getenv("ATT_TEXT_CACHE", "/index/att_text"))
ATT_TEXT_CACHE.mkdir(parents=True, exist_ok=True)

DATA_ATT_ROOT = Path(settings.data_root) / "attachments"        # /data/attachments
ADO_ATT_ROOT  = (Path(settings.ado_root) / "attachments") if settings.ado_root else None  # /ado/attachments

_STOP = False
log = logging.getLogger("attachments_ingest")

def _setup_logging() -> str:
    log_path = os.getenv("ATT_LOG_PATH", "/index/attachments_ingest.log")
    os.makedirs(os.path.dirname(log_path), exist_ok=True)
    logging.basicConfig(
        level=logging.DEBUG,
        format="%(asctime)s %(levelname)-7s %(name)s %(message)s",
        handlers=[logging.StreamHandler(), logging.FileHandler(log_path, mode="a", encoding="utf-8")],
    )
    log.info(f"[ATT] logfile={log_path}")
    return log_path

def _sig_handler(signum, frame):
    global _STOP
    _STOP = True
    log.info(f"[ATT] signal {signum} received; will stop after current file...")

for _sig in (signal.SIGINT, signal.SIGTERM):
    signal.signal(_sig, _sig_handler)

def _run_cmd(cmd: List[str], timeout: int) -> Tuple[int, bytes, bytes]:
    try:
        proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=timeout, check=False)
        return proc.returncode, proc.stdout, proc.stderr
    except subprocess.TimeoutExpired as e:
        return 124, b"", f"timeout: {e}".encode()

def _sniff_mime(path: Path) -> str:
    try:
        import magic
        return (magic.from_file(str(path), mime=True) or "").lower()
    except Exception:
        return ""

def _extract_txt(path: Path) -> str:
    try:
        raw = path.read_bytes()
        try:
            return raw.decode("utf-8")
        except UnicodeDecodeError:
            return raw.decode("latin-1", "ignore")
    except Exception:
        return ""

def _extract_image_text(path: Path) -> str:
    if not (OCR_ENABLED and OCR_IMAGES):
        return ""
    try:
        from PIL import Image
        import pytesseract
        im = Image.open(str(path))
        return pytesseract.image_to_string(im, lang=IMG_OCR_LANG, config=f"--psm 6 --oem 1 --dpi {OCR_DPI}") or ""
    except Exception:
        return ""

def _pdftotext_has_text(pdf_path: Path) -> Optional[str]:
    if not shutil.which("pdftotext"):
        return None
    rc, out, err = _run_cmd(["pdftotext", "-layout", str(pdf_path), "-"], PDF_TXT_TIMEOUT)
    if rc == 0 and out:
        return out.decode("utf-8", "ignore")
    return ""

def _ocr_pdf_to_text(pdf_path: Path) -> str:
    if not (OCR_ENABLED and OCR_PDF_IF_EMPTY):
        return ""
    if not (shutil.which("ocrmypdf") and shutil.which("pdftotext")):
        return ""
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as tmp:
        try:
            shutil.copyfile(str(pdf_path), tmp.name)
        except Exception:
            return ""
        args = [
            "ocrmypdf", "--force-ocr", "--skip-text",
            "--render-dpi", str(OCR_DPI),
            "--output-type", "pdf",
            "--optimize", "0",
            tmp.name, tmp.name
        ]
        if OCR_JOBS > 0:
            args.insert(1, "--jobs"); args.insert(2, str(OCR_JOBS))
        rc, out, err = _run_cmd(args, OCR_TIMEOUT_SEC)
        if rc != 0:
            log.debug(f"[ATT][OCR] ocrmypdf rc={rc} err={err.decode('utf-8','ignore')[:200]}")
            return ""
        rc2, out2, err2 = _run_cmd(["pdftotext", "-layout", tmp.name, "-"], PDF_TXT_TIMEOUT)
        if rc2 == 0 and out2:
            return out2.decode("utf-8", "ignore")
    return ""

def _extract_docx_text(path: Path) -> str:
    try:
        import docx
        doc = docx.Document(str(path))
        parts: List[str] = []
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)
        for tbl in doc.tables:
            for row in tbl.rows:
                parts.append("\t".join((c.text or "").strip() for c in row.cells))
        return "\n".join(parts)
    except Exception:
        return ""

def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts: List[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"\n--- SLIDE {i} ---\n")
            for shp in slide.shapes:
                if hasattr(shp, "text"):
                    t = (shp.text or "").strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)
    except Exception:
        return ""

def _extract_xlsx_text(path: Path, max_cells: int = 5000) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts: List[str] = []
        cells = 0
        for ws in wb.worksheets:
            parts.append(f"\n--- SHEET {ws.title} ---\n")
            for row in ws.iter_rows(values_only=True):
                line = ["" if v is None else str(v) for v in row]
                if any(line):
                    parts.append("\t".join(line))
                    cells += 1
                    if cells >= max_cells:
                        break
            if cells >= max_cells:
                break
        return "\n".join(parts)
    except Exception:
        return ""

def _truncate(s: str, max_chars: int) -> str:
    if not s: return s
    return s if len(s) <= max_chars else s[:max_chars] + "\n\n... [truncated]\n"

def _extract_text_for_file(path: Path) -> str:
    deny_ext = {
        ".rpm",".deb",".msi",".exe",".dll",".so",".whl",".jar",".iso",".img",
        ".tar",".tgz",".tar.gz",".tar.xz",".tar.bz2",".gz",".xz",".bz2",".7z",".zip",
        ".obj",".a",".lib"
    }
    ext  = path.suffix.lower()
    if ext in deny_ext:
        return ""
    mime = _sniff_mime(path)
    if mime.startswith("text/"):
        return _extract_txt(path)
    if ext == ".pdf" or "pdf" in mime:
        txt = _pdftotext_has_text(path)
        if txt:
            return txt
        return _ocr_pdf_to_text(path)
    if ext == ".docx":
        return _extract_docx_text(path)
    if ext == ".pptx":
        return _extract_pptx_text(path)
    if ext == ".xlsx":
        return _extract_xlsx_text(path)
    if ext in (".txt", ".csv", ".md", ".log", ".cfg", ".ini", ".conf", ".json", ".xml", ".yaml", ".yml",
               ".py", ".sh", ".rb", ".php", ".js", ".ts", ".java", ".go", ".cpp", ".cc", ".c", ".h", ".hpp",
               ".sql", ".ps1", ".bat"):
        return _extract_txt(path)
    if ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp", ".gif") or mime.startswith("image/"):
        return _extract_image_text(path)
    allow_mime = {
        "application/msword", "application/rtf", "application/json", "application/xml",
        "application/javascript", "application/x-javascript",
        "application/x-sh", "application/x-shellscript",
        "application/x-python", "application/x-perl", "application/x-ruby", "application/x-php",
        "application/vnd.ms-powerpoint", "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    }
    if mime.startswith("application/"):
        if mime in allow_mime:
            return _extract_txt(path)
        return ""
    return _extract_txt(path)

def _count_attachments(conn: sqlite3.Connection, spaces: Optional[List[str]]) -> int:
    params: List[Any] = []
    where = ""
    if spaces:
        marks = ",".join("?" for _ in spaces)
        where = f"WHERE d.space IN ({marks})"
        params.extend(spaces)
    sql = f"SELECT COUNT(*) FROM attachments a JOIN docs d ON d.id=a.content_id {where}"
    return conn.execute(sql, params).fetchone()[0]

def _iter_attachments(conn: sqlite3.Connection, spaces: Optional[List[str]], page: int = 500):
    params: List[Any] = []
    where = ""
    if spaces:
        marks = ",".join("?" for _ in spaces)
        where = f"WHERE d.space IN ({marks})"
        params.extend(spaces)
    base = f"""
      SELECT a.content_id, d.space, a.name, a.relpath, COALESCE(a.size,0)
      FROM attachments a JOIN docs d ON d.id = a.content_id
      {where}
      ORDER BY d.space, a.content_id, a.name
      LIMIT ? OFFSET ?
    """
    off = 0
    while True:
        rows = conn.execute(base, [*params, page, off]).fetchall()
        if not rows:
            break
        for r in rows:
            yield (r[0], r[1], r[2], r[3], r[4])
        off += page

def _resolve_abs_path(relpath: str) -> Path:
    """
    Map DB relpath to actual filesystem:
      - "ADO/<file>"      ->  /ado/attachments/<file>          (settings.ado_root)
      - everything else   ->  /data/attachments/<relpath>      (settings.data_root)
    """
    if relpath.startswith("ADO/"):
        if ADO_ATT_ROOT is None:
            return Path("/nonexistent")  # will miss
        return ADO_ATT_ROOT / relpath[len("ADO/"):]
    return DATA_ATT_ROOT / relpath

def _embed_and_upsert(store: NumpyStore, embedder, items: List[Dict[str, Any]]) -> int:
    if not items:
        return 0
    texts = [it["text"] for it in items]
    embs  = embedder.embed_texts(texts)
    return store.upsert_batch(items, embs)

def ingest(spaces_arg: str = "ALL", batch: int = 500) -> int:
    _setup_logging()
    conn = connect(settings.index_path)
    store = NumpyStore(VSConfig())
    embedder = get_embedder(dim_hint=int(os.getenv("EMBED_DIM", "1536")))

    spaces = None if spaces_arg == "ALL" else [s.strip() for s in spaces_arg.split(",") if s.strip()]
    total_up = 0
    scanned  = 0
    total_att = _count_attachments(conn, spaces)

    eff_flush_every = FLUSH_EVERY if FLUSH_EVERY > 0 else batch
    log.info(f"[ATT] index={settings.index_path} spaces={spaces_arg} batch={batch} "
             f"attachments_total={total_att} chunk_min_chars={CHUNK_MIN_CHARS} "
             f"OCR_ENABLED={int(OCR_ENABLED)} OCR_PDF_IF_EMPTY={int(OCR_PDF_IF_EMPTY)} "
             f"OCR_IMAGES={int(OCR_IMAGES)} OCR_DPI={OCR_DPI} OCR_JOBS={OCR_JOBS} "
             f"FLUSH_EVERY={eff_flush_every}")

    try:
        for content_id, space, name, relpath, size in _iter_attachments(conn, spaces, page=batch):
            if _STOP:
                log.info(f"[ATT] stop requested; exit after scanned={scanned} up_total={total_up}")
                break

            abs_path = _resolve_abs_path(relpath)
            if not abs_path.is_file():
                log.info(f"[ATT][MISS] {space} {relpath} (file not found)")
                scanned += 1
                if scanned % eff_flush_every == 0:
                    store.flush()
                continue

            key = hashlib.sha256(f"{abs_path}:{abs_path.stat().st_size}:{int(abs_path.stat().st_mtime)}".encode()).hexdigest()
            cache_file = ATT_TEXT_CACHE / (key + ".txt")

            if cache_file.is_file():
                try:
                    text = cache_file.read_text(encoding="utf-8", errors="ignore")
                except Exception:
                    text = ""
                src = "cache"
            else:
                text = _truncate(_extract_text_for_file(abs_path), ATT_MAX_CHARS)
                try:
                    cache_file.write_text(text or "", encoding="utf-8")
                except Exception:
                    pass
                src = "extract"

            if not text or len(text) < CHUNK_MIN_CHARS:
                log.info(f"[ATT][SKIP] {space} {relpath} name={name} reason=short({len(text) if text else 0}) src={src}")
                scanned += 1
                if scanned % eff_flush_every == 0:
                    store.flush()
                continue

            doc_id = f"ATT:{content_id}:{relpath}"
            base_meta = {
                "id": doc_id, "space": space, "type": "attachment",
                "title": name, "url": f"/attachments/{relpath}", "updated_at": ""
            }
            items: List[Dict[str, Any]] = []
            for ch in iter_chunks({"id": doc_id, "space": space, "type": "attachment", "title": name, "text": text}):
                meta = dict(base_meta)
                meta["chunk_ix"] = ch["chunk_ix"]
                meta["text"] = ch["text"]
                items.append(meta)

            n = _embed_and_upsert(store, embedder, items)
            total_up += n
            scanned  += 1
            log.info(f"[ATT][OK] {space} {relpath} name={name} src={src} chunks={n} scanned={scanned} up_total={total_up}")

            if scanned % eff_flush_every == 0:
                store.flush()

    finally:
        store.flush()
        log.info(f"[ATT] DONE: total_chunks_upserted={total_up} scanned={scanned}")

    return total_up

def main():
    import argparse
    ap = argparse.ArgumentParser(description="Extract + embed attachments.")
    ap.add_argument("--spaces", default=os.getenv("ATT_SPACES", "ALL"),
                    help='ALL or comma-separated list, e.g. "OTT,ADO"')
    ap.add_argument("--batch", type=int, default=int(os.getenv("ATT_BATCH", "500")))
    args = ap.parse_args()
    ingest(args.spaces, args.batch)

if __name__ == "__main__":
    main()
